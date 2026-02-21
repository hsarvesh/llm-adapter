"""Parser for document formats: PDF, DOCX, PPTX."""

import io
import structlog

from parsers.base import BaseParser

logger = structlog.get_logger(__name__)


class DocumentParser(BaseParser):
    """Handles PDF, DOCX, and PPTX documents."""

    @property
    def supported_extensions(self) -> list[str]:
        return [".pdf", ".docx", ".pptx"]

    def parse(self, file_bytes: bytes, filename: str) -> str:
        ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""

        if ext == ".pdf":
            return self._parse_pdf(file_bytes, filename)
        elif ext == ".docx":
            return self._parse_docx(file_bytes, filename)
        elif ext == ".pptx":
            return self._parse_pptx(file_bytes, filename)
        else:
            return f"[Unsupported document format: {ext}]"

    def _parse_pdf(self, file_bytes: bytes, filename: str) -> str:
        """Extract text from PDF using pdfplumber (primary) with PyMuPDF fallback."""
        # Try pdfplumber first (better table extraction)
        try:
            import pdfplumber

            pages_text = []
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""

                    # Also try to extract tables
                    tables = page.extract_tables()
                    table_text = ""
                    for table in tables:
                        if table:
                            for row in table:
                                row_str = " | ".join(str(cell or "") for cell in row)
                                table_text += row_str + "\n"
                            table_text += "\n"

                    page_content = text
                    if table_text:
                        page_content += "\n\n[Tables]\n" + table_text

                    if page_content.strip():
                        pages_text.append(f"--- Page {i + 1} ---\n{page_content.strip()}")

            if pages_text:
                result = "\n\n".join(pages_text)
                logger.info("pdf_parsed_pdfplumber", filename=filename, pages=len(pages_text))
                return result
        except Exception as e:
            logger.warning("pdfplumber_fallback", filename=filename, error=str(e))

        # Fallback to PyMuPDF
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(stream=file_bytes, filetype="pdf")
            pages_text = []
            for i, page in enumerate(doc):
                text = page.get_text().strip()
                if text:
                    pages_text.append(f"--- Page {i + 1} ---\n{text}")
            doc.close()

            result = "\n\n".join(pages_text)
            logger.info("pdf_parsed_pymupdf", filename=filename, pages=len(pages_text))
            return result if result else "(No text content found in PDF)"
        except Exception as e:
            logger.error("pdf_parse_error", filename=filename, error=str(e))
            return f"[Error extracting PDF content: {str(e)}]"

    def _parse_docx(self, file_bytes: bytes, filename: str) -> str:
        """Extract text from DOCX using python-docx."""
        try:
            from docx import Document

            doc = Document(io.BytesIO(file_bytes))
            paragraphs = []

            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # Preserve heading levels
                    if para.style and para.style.name.startswith("Heading"):
                        level = para.style.name.replace("Heading ", "").strip()
                        try:
                            prefix = "#" * int(level)
                            text = f"{prefix} {text}"
                        except ValueError:
                            pass
                    paragraphs.append(text)

            # Extract tables
            for table in doc.tables:
                table_rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_rows.append(" | ".join(cells))
                if table_rows:
                    paragraphs.append("\n[Table]\n" + "\n".join(table_rows))

            result = "\n\n".join(paragraphs)
            logger.info("docx_parsed", filename=filename, paragraphs=len(paragraphs))
            return result if result else "(No text content found in DOCX)"
        except Exception as e:
            logger.error("docx_parse_error", filename=filename, error=str(e))
            return f"[Error extracting DOCX content: {str(e)}]"

    def _parse_pptx(self, file_bytes: bytes, filename: str) -> str:
        """Extract text from PPTX using python-pptx."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(file_bytes))
            slides_text = []

            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                texts.append(text)
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            texts.append(" | ".join(cells))

                if texts:
                    slides_text.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))

            result = "\n\n".join(slides_text)
            logger.info("pptx_parsed", filename=filename, slides=len(slides_text))
            return result if result else "(No text content found in PPTX)"
        except Exception as e:
            logger.error("pptx_parse_error", filename=filename, error=str(e))
            return f"[Error extracting PPTX content: {str(e)}]"
